"""BCFF WP スライド生成スクリプト v4
場所: n2p-wp-design skill/generate_slides.py
"""
import os
import copy
import math
import re
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

BASE   = RGBColor(0xFC, 0xFC, 0xFC)
MAIN1  = RGBColor(0x11, 0x11, 0x11)
MAIN2  = RGBColor(0x43, 0x43, 0x43)
ACCENT = RGBColor(0x16, 0x65, 0xF2)
LINE   = RGBColor(0xE6, 0xEB, 0xED)
CR     = RGBColor(0x92, 0x96, 0x99)
BAR    = RGBColor(0x93, 0x96, 0x99)

SKILL_DIR  = r"C:\Users\kazug\.claude\N2Pホワイトペーパー\.claude\skills\n2p-wp-design skill"
LOGO_W     = SKILL_DIR + r"\assets\symbol_white.png"
LOGO_B     = SKILL_DIR + r"\assets\symbol_black.png"
INSERT_DIR = SKILL_DIR + r"\inserts"
OUTPUT_DIR = r"C:\Users\kazug\.claude\N2Pホワイトペーパー\BCFF WP"



# ─── ユーティリティ ──────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def remove_shadow(shape):
    """図形のドロップシャドウを完全に無効化"""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        for el in list(spPr):
            if el.tag in (qn('a:effectLst'), qn('a:effectDag')):
                spPr.remove(el)
        etree.SubElement(spPr, qn('a:effectLst'))
    style = sp.find(qn('p:style'))
    if style is not None:
        eff = style.find(qn('a:effectRef'))
        if eff is not None:
            eff.set('idx', '0')
            for c in list(eff):
                eff.remove(c)
            etree.SubElement(eff, qn('a:srgbClr')).set('val', '000000')

def add_rect(slide, l, t, w, h, fill_color, no_line=True):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Cm(l), Cm(t), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if no_line:
        shp.line.fill.background()
    remove_shadow(shp)
    return shp

def parse_marked(text):
    """==text== 形式を解析し [(text, is_highlighted), ...] を返す"""
    parts = re.split(r'==(.+?)==', text)
    result = []
    for i, part in enumerate(parts):
        if part:
            result.append((part, i % 2 == 1))
    return result

def apply_highlight(run, color_hex="F4F400"):
    """ランの <a:rPr> に <a:highlight> を追加してラインマーカーを適用"""
    r_elem = run._r
    rPr = r_elem.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.Element(qn('a:rPr'))
        r_elem.insert(0, rPr)
    for el in list(rPr):
        if el.tag == qn('a:highlight'):
            rPr.remove(el)
    highlight = etree.SubElement(rPr, qn('a:highlight'))
    etree.SubElement(highlight, qn('a:srgbClr')).set('val', color_hex)

def _run_with_marks(p, text, size, bold=False, color=MAIN2):
    """==text== マーカーを解析し、ハイライト付きランを段落に追加する"""
    for content, highlighted in parse_marked(text):
        r = p.add_run()
        r.text = content
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if highlighted:
            apply_highlight(r)

def _set_para_fmt(p, lsp=150):
    """行間・前後スペース0 を設定"""
    pPr = p._p.get_or_add_pPr()
    for tag in (qn('a:lnSpc'), qn('a:spcBef'), qn('a:spcAft')):
        for el in pPr.findall(tag):
            pPr.remove(el)
    lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
    etree.SubElement(lnSpc, qn('a:spcPct')).set('val', str(lsp * 1000))
    spcB = etree.SubElement(pPr, qn('a:spcBef'))
    etree.SubElement(spcB, qn('a:spcPts')).set('val', '0')
    spcA = etree.SubElement(pPr, qn('a:spcAft'))
    etree.SubElement(spcA, qn('a:spcPts')).set('val', '0')

def _run(p, text, size, bold=False, color=MAIN2):
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color

def add_simple_tb(slide, l, t, w, h, text,
                  size=12, bold=False, color=MAIN1,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.TOP,
                  pl=0.25, pt_=0.25, pr=0.25, pb=0.25, wrap=True):
    """単一段落テキストボックス"""
    tx = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tx.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = v_anchor
    tf.margin_left = Cm(pl);  tf.margin_right  = Cm(pr)
    tf.margin_top  = Cm(pt_); tf.margin_bottom = Cm(pb)
    p = tf.paragraphs[0]
    p.alignment = align
    _set_para_fmt(p, lsp=100)
    _run(p, text, size, bold, color)
    return tx

def add_body_tb(slide, l, t, w, h, text,
                size=13, color=MAIN2,
                pl=0, pr=0, pt_=0.5, pb=0.5):
    """本文テキストボックス（1.5行間・空行除去）"""
    lines = [ln for ln in text.split('\n') if ln.strip()]
    tx = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Cm(pl);  tf.margin_right  = Cm(pr)
    tf.margin_top  = Cm(pt_); tf.margin_bottom = Cm(pb)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        _set_para_fmt(p, lsp=150)
        _run_with_marks(p, line, size, False, color)
    return tx

def add_headed_tb(slide, l, t, w, h, heading, body_text,
                  hsize=13, bsize=11,
                  pl=1.0, pr=1.0, pt_=0.25, pb=0.25):
    """見出し＋本文（1.5行間）テキストボックス"""
    lines = [ln for ln in body_text.split('\n') if ln.strip()]
    tx = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Cm(pl);  tf.margin_right  = Cm(pr)
    tf.margin_top  = Cm(pt_); tf.margin_bottom = Cm(pb)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p0, lsp=150)
    _run(p0, heading, hsize, bold=True, color=ACCENT)
    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        _set_para_fmt(p, lsp=150)
        _run_with_marks(p, line, bsize, False, MAIN2)
    return tx

def add_combined_case_tb(slide, l, t, w, h, overview, flow,
                          hsize=13, bsize=11,
                          pl=1.0, pr=1.0, pt_=0.25, pb=0.25):
    """キャンペーン概要＋7pt空白行＋キャンペーンフロー を1つのTBに統合（1.5行間）"""
    tx = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Cm(pl);  tf.margin_right  = Cm(pr)
    tf.margin_top  = Cm(pt_); tf.margin_bottom = Cm(pb)

    def para(first=False):
        return tf.paragraphs[0] if first else tf.add_paragraph()

    # 「キャンペーン概要」見出し（アクセントカラー）
    p = para(first=True)
    p.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p, lsp=150)
    _run(p, "キャンペーン概要", hsize, bold=True, color=ACCENT)

    # 概要本文（1段落・自然折り返し、==mark== ハイライト対応）
    p = para()
    p.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p, lsp=150)
    _run_with_marks(p, overview, bsize, False, MAIN2)

    # 7pt 空白行（スペース文字で行高を確保）
    p = para()
    p.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p, lsp=150)
    _run(p, " ", 7, False, MAIN2)

    # 「キャンペーンフロー」見出し（アクセントカラー）
    p = para()
    p.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p, lsp=150)
    _run(p, "キャンペーンフロー", hsize, bold=True, color=ACCENT)

    # フロー各ステップ（1ステップ＝1段落、==mark== ハイライト対応）
    for line in [ln for ln in flow.split('\n') if ln.strip()]:
        p = para()
        p.alignment = PP_ALIGN.LEFT
        _set_para_fmt(p, lsp=150)
        _run_with_marks(p, line, bsize, False, MAIN2)

    return tx

def est_h_combined(overview, flow, w_cm,
                   hsize=13, bsize=11, lsp=1.5,
                   ml=1.0, mr=1.0, mt=0.25, mb=0.25, buf=1.12):
    """概要＋スペーサー＋フロー 統合TBの高さを推定（cm）"""
    avail_pt = (w_cm - ml - mr) * 28.35
    cpline = max(1.0, avail_pt / (bsize * 0.85))

    def lh(pt):
        return pt * lsp * 2.54 / 72

    h = mt + mb
    h += lh(hsize)                                              # 概要見出し
    h += math.ceil(max(1, len(overview)) / cpline) * lh(bsize) # 概要本文
    h += lh(7)                                                  # 7pt空白行
    h += lh(hsize)                                              # フロー見出し
    for line in [ln for ln in flow.split('\n') if ln.strip()]:
        h += math.ceil(max(1, len(line)) / cpline) * lh(bsize) # フロー各行
    return h * buf

def est_h(heading, body_text, hsize, bsize, lsp, w_cm,
          ml=1.0, mr=1.0, mt=0.25, mb=0.25, buf=1.12):
    """見出し+本文TBの高さを推定（cm）"""
    avail_pt = (w_cm - ml - mr) * 28.35
    cpline = max(1.0, avail_pt / (bsize * 0.85))

    def lh(pt):
        return pt * lsp * 2.54 / 72

    h = mt + mb + lh(hsize)
    for line in [ln for ln in body_text.split('\n') if ln.strip()]:
        h += math.ceil(max(1, len(line)) / cpline) * lh(bsize)
    return h * buf

def add_logo(slide, white=True):
    path = LOGO_W if white else LOGO_B
    if os.path.exists(path):
        slide.shapes.add_picture(path, Cm(23.13), Cm(0.82), Cm(1.72), Cm(0.57))

def add_footer(slide, pn, on_accent=False):
    lc = BASE if on_accent else LINE
    tc = BASE if on_accent else CR
    add_rect(slide, 0.59, 13.58, 24.23, 0.03, lc)
    add_simple_tb(slide, 22.73, 13.58, 2.12, 0.71, str(pn),
                  size=10, color=tc, align=PP_ALIGN.RIGHT,
                  v_anchor=MSO_ANCHOR.MIDDLE)
    add_simple_tb(slide, 0.75, 13.52, 7.55, 0.81, "©NONAME Produce Inc.",
                  size=7, color=tc, v_anchor=MSO_ANCHOR.MIDDLE)


# ─── スライドビルダー ──────────────────────────────────

def build_cover(prs, pn):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_footer(sl, pn)

def build_toc(prs, chapters, pn):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_rect(sl, 12.5, 0, 12.95, 14.29, ACCENT)
    add_logo(sl, white=True)
    add_rect(sl, 0.59, 13.58, 24.23, 0.03, LINE)
    add_simple_tb(sl, 0.75, 13.52, 7.55, 0.81, "©NONAME Produce Inc.",
                  size=7, color=CR, v_anchor=MSO_ANCHOR.MIDDLE)
    add_simple_tb(sl, 22.73, 13.58, 2.12, 0.71, str(pn),
                  size=10, color=BASE, align=PP_ALIGN.RIGHT, v_anchor=MSO_ANCHOR.MIDDLE)
    add_simple_tb(sl, 1, 5.87, 1.72, 1.11, "00.",
                  size=14, color=MAIN1, v_anchor=MSO_ANCHOR.MIDDLE)
    add_simple_tb(sl, 1.79, 5.83, 6.13, 1.54, "目次",
                  size=24, bold=True, color=MAIN1, v_anchor=MSO_ANCHOR.MIDDLE)

    # 章リスト TB（幅3.84in × 高4.5in、X=5.27in、スライド上下中央）
    tb_w = Inches(3.84)
    tb_h = Inches(4.5)
    tb_x = Inches(5.27)
    tb_y = Inches((14.29 / 2.54 - 4.5) / 2)

    lines = []
    for i, ch in enumerate(chapters):
        lines.append(f"{str(i+1).zfill(2)}.{ch}")
        if i < len(chapters) - 1:
            lines.append("")

    tx = sl.shapes.add_textbox(tb_x, tb_y, tb_w, tb_h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.25); tf.margin_right  = Cm(0.25)
    tf.margin_top  = Cm(0.25); tf.margin_bottom = Cm(0.25)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line.strip():
            _set_para_fmt(p, lsp=115)
            _run(p, line, 14, bold=True, color=BASE)
        else:
            _set_para_fmt(p, lsp=50)
            _run(p, "", 7, color=BASE)

def build_chapter(prs, ch_num, title, pn):
    """章題ページ。タイトルは左上揃え。"""
    sl = blank_slide(prs)
    set_bg(sl, ACCENT)
    add_logo(sl, white=True)
    add_footer(sl, pn, on_accent=True)
    add_simple_tb(sl, 1, 5.87, 1.72, 1.11, f"{str(ch_num).zfill(2)}.",
                  size=14, color=BASE, v_anchor=MSO_ANCHOR.MIDDLE)
    # タイトル：左揃え・上揃え（spec: 左上）
    add_simple_tb(sl, 1.79, 5.83, 19.68, 2.5, title,
                  size=24, bold=True, color=BASE,
                  v_anchor=MSO_ANCHOR.TOP, wrap=True,
                  pl=0.25, pt_=0.25, pr=0.25, pb=0.25)

def build_general(prs, title, body, pn):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_rect(sl, 0, 1.08, 1, 0.03, BAR)
    add_logo(sl, white=False)
    add_footer(sl, pn)
    add_simple_tb(sl, 1, 0.69, 21.58, 0.81, title,
                  size=19, bold=True, color=ACCENT,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)
    add_body_tb(sl, 1, 1.51, 23.4, 11.5, body, size=13, color=MAIN2)

def build_case1(prs, case_num, name, company, overview, flow, pn):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_rect(sl, 0, 0, 25.4, 3.01, ACCENT)
    add_rect(sl, 0, 1.08, 1, 0.03, BASE)
    add_logo(sl, white=True)
    add_footer(sl, pn)

    add_simple_tb(sl, 1, 0.69, 21.58, 0.81,
                  f"{str(case_num).zfill(2)}.{name}",
                  size=19, bold=True, color=BASE,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)
    add_simple_tb(sl, 1, 1.76, 21.58, 0.56, f"by {company}",
                  size=13, color=BASE,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)

    tb_w = 12.7
    tb_h = est_h_combined(overview, flow, tb_w)
    add_combined_case_tb(sl, 0, 3.3, tb_w, tb_h, overview, flow)

def build_case2(prs, case_num, name, company, point, pn):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_rect(sl, 0, 0, 25.4, 3.01, ACCENT)
    add_rect(sl, 0, 1.08, 1, 0.03, BASE)
    add_logo(sl, white=True)
    add_footer(sl, pn)

    add_simple_tb(sl, 1, 0.69, 21.58, 0.81,
                  f"{str(case_num).zfill(2)}.{name}",
                  size=19, bold=True, color=BASE,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)
    add_simple_tb(sl, 1, 1.76, 21.58, 0.56, f"by {company}",
                  size=13, color=BASE,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)

    pt_h = est_h("キャンペーンのポイント", point, 13, 11, 1.5, 12.7)
    add_headed_tb(sl, 0, 3.3, 12.7, pt_h,
                  "キャンペーンのポイント", point, hsize=13, bsize=11)


# ─── PPTX 差し込み（PowerPoint COM を使用）────────────────

def _com_insert_service_slides(out_path, service_map, ch_positions):
    """
    PowerPoint COM を使ってサービスファイルを章タイトルの直後に差し込む。
    テーマ・フォント・カラーをソース通りに保持する唯一の確実な方法。

    service_map:  {chapter_num: src_file_path}
    ch_positions: {chapter_num: 1-indexed slide number in target}
    """
    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True   # False だと一部操作が失敗するため True で起動

    skipped = []
    try:
        tgt = ppt.Presentations.Open(
            os.path.abspath(out_path),
            False,   # ReadOnly
            False,   # Untitled
            False)   # WithWindow（バックグラウンド）

        offset = 0  # 先の挿入でスライド番号がずれる分を補正

        for ch in sorted(ch_positions.keys()):
            src = service_map.get(ch, "")
            insert_after = ch_positions[ch] + offset

            if not src or not os.path.exists(src):
                skipped.append(f"第{ch}章（{os.path.basename(src or '未設定')} が見つかりません）")
                continue

            abs_src = os.path.abspath(src)

            # スライド枚数を確認
            tmp = ppt.Presentations.Open(abs_src, True, False, False)
            n = tmp.Slides.Count
            tmp.Close()

            # 差し込み（ソース書式を保持）
            tgt.Slides.InsertFromFile(abs_src, insert_after)
            offset += n
            print(f"  INSERT: 第{ch}章 → {n}枚 after slide {insert_after}")

        tgt.Save()
        tgt.Close()

    finally:
        ppt.Quit()
        pythoncom.CoUninitialize()

    return skipped


# ─── メイン ────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Cm(25.4)
    prs.slide_height = Cm(14.29)
    pn = 1
    skipped = []

    chapters = [
        "なぜ従来のキャンペーンは限界を迎えたのか？",
        "発話を生み出す「マイレージキャンペーン」の仕組み",
        "成功事例に学ぶ発話設計",
        "これらを最短で実現する「BirdCall for Fun」",
        "まとめ",
        "BirdCall for Funのご案内",
        "その他のご案内",
    ]

    build_cover(prs, pn); pn += 1
    build_toc(prs, chapters, pn); pn += 1

    # ── 第1章 ──
    build_chapter(prs, 1, "なぜ従来のキャンペーンは限界を迎えたのか？", pn); pn += 1

    build_general(prs, "課題は露出量より「誰が語るか」",
        "現代の情報環境では「誰が発信するか」が信頼度を大きく左右する。\n"
        "企業からの発信は警戒・回避されやすく、ユーザーからの発信は信頼・共感されやすい。\n"
        "広告の課題は「露出量の不足」ではなく、「発信主体（誰が語るか）」にある。\n"
        "UGCを活用したキャンペーンは非UGCと比較してコンバージョンが29%高い（Hootsuite 2025）",
        pn); pn += 1

    build_general(prs, "参加はされるが、何も残らない",
        "フォロー＆リポストに代表される従来型キャンペーンは、ユーザーの「瞬間的な参加」しか生まない。\n"
        "ボタン1つで完結する行動にはブランドへの理解や愛着が伴わず、キャンペーン終了後にはユーザーとの関係が何も残らない。\n"
        "さらに懸賞応募目的のアカウントが多く集まり、本来獲得したい層との接点にならないという課題もある。",
        pn); pn += 1

    build_general(prs, "これからのプロモーションに必要なのは「発話の質」",
        "参加数ではなく、ユーザーが「自分の言葉で語る」発話の質こそが重要な時代になった。\n"
        "体験を伴った投稿・引用コメント・画像付きのレビューなどは、SNSアルゴリズムにもユーザーにも評価されやすく、自然な拡散を生む。\n"
        "では、この「質の高い発話」はどうすれば設計できるのか？",
        pn); pn += 1

    # ── 第2章 ──
    build_chapter(prs, 2, "発話を生み出す「マイレージキャンペーン」の仕組み", pn); pn += 1

    build_general(prs, "参加を「点」から「線」へ",
        "マイレージキャンペーンとは、ユーザーの行動をポイント化し、継続的な参加を促す設計の仕組み。\n"
        "従来の「1回の参加（点）」を「複数回の関与（線）」に変えることで、ブランドとの関係を積み上げていく。",
        pn); pn += 1

    build_general(prs, "発話は偶然ではなく、設計によって生まれる",
        "ユーザーの発話は自然に増えるのを待つものではなく、仕組みによって引き出せる。\n"
        "マイレージ設計では3つのメカニズムが働く：\n"
        "「投稿するとポイントが貯まる（動機の明確化）」\n"
        "「行動量が可視化される（自己強化）」\n"
        "「他ユーザーと比較できる（競争意識）」",
        pn); pn += 1

    build_general(prs, "質の高い発話が、アルゴリズムを動かす",
        "マイレージ型ではユーザーが「より高いポイントを得るために行動を工夫」するため、引用投稿・コメント・画像付き投稿といった質の高いUGCが自然と増える。\n"
        "これらはXのアルゴリズム高評価指標と一致している：\n"
        "返信 ×13.5　リポスト ×20　いいね ×1\n"
        "広告費なしの自然な拡散につながる。",
        pn); pn += 1

    build_general(prs, "習慣的な参加を生む「蓄積・競争・達成」の設計",
        "継続的な発話を生み出すのは3つの心理設計：\n"
        "「蓄積」──やればやるほど得をする\n"
        "「競争」──ランキングで他者と比較できる\n"
        "「達成」──報酬・称号で満足感を得る\n"
        "この3つが組み合わさることで、単発の参加が習慣的な行動に変わる。",
        pn); pn += 1

    # ── 第3章 ──
    build_chapter(prs, 3, "成功事例に学ぶ発話設計", pn); pn += 1

    build_case1(prs, 1, "ブレンユー1周年記念マイレージキャンペーン", "ゼブラ株式会社",
        # overview: 改行なし、自然に折り返す
        overview="ゼブラ「ブレンユー」の1周年を記念して実施したマイレージキャンペーン。ハッシュタグ投稿（50マイル）をいいね（10マイル）より高く設定し、ユーザーが自発的にUGCを生み出す設計を採用。継続的なアクションを促す構造により、実施期間全体を通じてエンゲージメントが持続した。",
        flow=(
            "1. 公式アカウント（@blen_u）をフォロー\n"
            "2. BCfF LPにアクセスしてアカウント連携\n"
            "3. 公式アカウントの投稿にリアクションでマイル獲得\n"
            "4. 指定ハッシュタグ付き投稿でさらにマイル獲得（50マイル）\n"
            "5. 貯めたマイルでプレゼントに応募"
        ),
        pn=pn); pn += 1

    build_case2(prs, 1, "ブレンユー1周年記念マイレージキャンペーン", "ゼブラ株式会社",
        point="「行動の価値に差をつける」というシンプルな設計変更だけで、ユーザーが積極的にUGCを生み出す動きが生まれた。マイレージ設計の基本原則（行動の細かい設計）が機能した典型例。",
        pn=pn); pn += 1

    build_case1(prs, 2, "コスメブランド30周年キャンペーン", "（企業名確認中）",
        overview="化粧品ブランドの周年記念として実施したUGC活性化型マイレージキャンペーン。「指定ハッシュタグ付きの一般ユーザー投稿へのリアクション」でもポイントが貯まる設計を採用。これにより、ファン同士の交流が促され、ブランドに関する投稿が連鎖的に広がった。",
        flow=(
            "1. 公式アカウントをフォロー\n"
            "2. BCfF LPにアクセスしてアカウント連携\n"
            "3. 公式アカウントへのリアクションでマイル獲得\n"
            "4. 指定ハッシュタグ付き一般ユーザー投稿へのリアクションでもマイル獲得\n"
            "5. 貯めたマイルでプレゼントに応募"
        ),
        pn=pn); pn += 1

    build_case2(prs, 2, "コスメブランド30周年キャンペーン", "（企業名確認中）",
        point="「一般ユーザーの投稿へのリアクション」をポイント対象にするという設計上の工夫が、ユーザー同士のインタラクションを誘発し、拡散の連鎖を生んだ。リポスト+引用リポストが約7万件に達したことが、コミュニティ形成の効果を示している。",
        pn=pn); pn += 1

    build_case1(prs, 3, "EBiDANマイルキャンペーン第1弾", "テレビ東京",
        overview="テレビ東京の番組「DAN!DAN!EBiDAN!」と連動したIPコラボ型マイルキャンペーン。公式アカウントへのリアクションに応じてマイルが加算され、貯まったマイルで豪華プレゼントへの応募や限定動画の視聴が可能な設計。ファンの熱量を報酬設計で引き出し、自発的なUGC創出につなげた。",
        flow=(
            "1. 公式アカウントをフォロー\n"
            "2. BCfF LPにアクセスしてアカウント連携\n"
            "3. 公式アカウントの投稿に指定のリアクションを送る\n"
            "4. マイルを獲得（いいね10・リポスト30・引用50・リプライ50・ハッシュタグ50マイル）\n"
            "5. 貯めたマイルで豪華プレゼントへの応募 & 限定動画を視聴"
        ),
        pn=pn); pn += 1

    build_case2(prs, 3, "EBiDANマイルキャンペーン第1弾", "テレビ東京",
        point="IPの持つファンの熱量を設計で引き出し、真のファンによるUGCを大量創出。広告費・TV告知なしでのトレンド入りは、有機的な拡散がいかに強力かを示している。",
        pn=pn); pn += 1

    build_general(prs, "成功事例から見えた、発話設計の3原則",
        "3つの事例に共通する成功パターン：\n"
        "「行動が細かく設計されている」──何をすればポイントが得られるか明確\n"
        "「継続する理由がある」──蓄積・競争・達成の仕組みが機能している\n"
        "「ユーザー同士の関係性が生まれている」──発話がコミュニティを形成しさらなる参加を促す",
        pn); pn += 1

    # ── 第4章 ──
    build_chapter(prs, 4, "これらを最短で実現する「BirdCall for Fun」", pn); pn += 1

    build_general(prs, "マイレージ設計には専門知識とコストが必要",
        "マイレージキャンペーンの効果は事例で示した通りだが、自社で設計・実装・運用するには高い専門性と工数が必要になる。\n"
        "・ポイント付与のロジック設計\n"
        "・ランキングシステムの構築\n"
        "・報酬の管理・配布\n"
        "・データ分析と最適化\n"
        "これらを全て自社でまかなうことは、多くのブランドにとって大きな負担となる。",
        pn); pn += 1

    build_general(prs, "「BirdCall for Fun」でその複雑さを解決する",
        "「BirdCall for Fun」はXのマイレージキャンペーンの設計・システム・運用管理を丸ごとパッケージ化したサービス。\n"
        "専門知識がなくてもすぐに実行でき、公式アカウントへのアクション（返信・引用・リポスト）だけでなく、ファン同士のハッシュタグ投稿へのリアクションもポイント対象にできる点が特徴。\n"
        "マイレージ設計の壁を取り除き、誰でも発話を設計できる環境を提供する。",
        pn); pn += 1

    build_general(prs, "3つの設計で発話を引き出す",
        "BirdCall for Funは以下の3つの設計でマイレージキャンペーンを実現する：\n"
        "①ポイント設計──公式アカウントへのアクション・ハッシュタグ投稿・ファン同士のインタラクションを自動スコアリング。1日の上限・各アクションのポイント数を自由に設定可能。\n"
        "②報酬設計──貯まったポイントは景品交換・後日抽選・インスタントウィンの3方式に対応。\n"
        "③連携設計──インスタントウィン・マストバイ等の他キャンペーンとのミッション形式の組み合わせも可能。",
        pn); pn += 1

    # ── 第5章 ──
    build_chapter(prs, 5, "まとめ", pn); pn += 1

    build_general(prs, "発話を設計する時代へ",
        "広告は「届ける」ものから「語られる」ものへと変化した。\n"
        "求められるのは、ユーザーの発話を偶発的に待つのではなく、構造的に生み出す設計である。\n"
        "その設計手法が「マイレージキャンペーン」であり、最短の実現手段が「BirdCall for Fun」である。",
        pn); pn += 1

    build_general(prs, "n2pと、発話が生まれるキャンペーンを設計しましょう",
        "n2pはSNSキャンペーンの設計・実施・分析を一気通貫でサポートするプロモーション制作会社。\n"
        "BirdCall for Funを活用したマイレージキャンペーンの設計から運用まで、まずはn2pに相談することが最初のステップ。",
        pn); pn += 1

    # ── 第6章（章タイトルのみ生成し、位置を記録）──
    build_chapter(prs, 6, "BirdCall for Funのご案内", pn); pn += 1
    ch6_pos = len(prs.slides)   # 章タイトルの 1-indexed スライド番号

    # ── 第7章（章タイトルのみ生成し、位置を記録）──
    build_chapter(prs, 7, "その他のご案内", pn); pn += 1
    ch7_pos = len(prs.slides)

    # ── Phase1: python-pptx 生成分を保存 ──
    out = os.path.join(OUTPUT_DIR, "BCFF-slides.pptx")
    prs.save(out)
    print(f"本文スライド生成完了: {len(prs.slides)}枚 -> {out}")

    # ── Phase2: PowerPoint COM でサービススライドを差し込み ──
    service_map = {
        6: os.path.join(INSERT_DIR, "n2p-wp-BCFF.pptx"),
        7: os.path.join(INSERT_DIR, "n2p-wp-other.pptx"),
    }
    com_skipped = _com_insert_service_slides(
        out, service_map, {6: ch6_pos, 7: ch7_pos})

    print("差し込み完了（PowerPoint COM）")
    for s in com_skipped:
        print(f"  SKIP: {s}")

if __name__ == "__main__":
    main()
