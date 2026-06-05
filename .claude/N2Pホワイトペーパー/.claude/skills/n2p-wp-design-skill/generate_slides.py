"""n2p WP 汎用スライド生成スクリプト v5
使い方: python generate_slides.py <*-wp-structure.md> [ACCENT_RRGGBB]
場所: n2p-wp-design skill/generate_slides.py
"""
import os
import sys
import copy
import math
import re
import tempfile
import unicodedata
import matplotlib
matplotlib.use('Agg')   # GUIなしで画像生成
import matplotlib.pyplot as plt
from matplotlib import font_manager
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

BASE   = RGBColor(0xFC, 0xFC, 0xFC)
MAIN1  = RGBColor(0x11, 0x11, 0x11)
MAIN2  = RGBColor(0x43, 0x43, 0x43)
ACCENT = RGBColor(0x16, 0x65, 0xF2)
LINE   = RGBColor(0xE6, 0xEB, 0xED)
CR     = RGBColor(0x92, 0x96, 0x99)
BAR    = RGBColor(0x93, 0x96, 0x99)

SKILL_DIR  = r"C:\Users\kazug\.claude\N2Pホワイトペーパー\.claude\skills\n2p-wp-design-skill"
LOGO_W     = SKILL_DIR + r"\assets\symbol_white.png"
LOGO_B     = SKILL_DIR + r"\assets\symbol_black.png"
INSERT_DIR = SKILL_DIR + r"\inserts"



# ─── ユーティリティ ──────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def remove_shadow(shape):
    """図形のドロップシャドウを完全に無効化"""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is not None:
        for el in list(spPr):
            if el.tag in (qn('a:effectLst'), qn('a:effectDag')):
                spPr.remove(el)
        etree.SubElement(spPr, qn('a:effectLst'))
    style = sp.find(qn('p:style'))
    if style is not None:
        eff = style.find(qn('a:effectRef'))
        if eff is not None:
            eff.set('idx', '0')
            for c in list(eff):
                eff.remove(c)
            etree.SubElement(eff, qn('a:srgbClr')).set('val', '000000')

def add_rect(slide, l, t, w, h, fill_color, no_line=True):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Cm(l), Cm(t), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if no_line:
        shp.line.fill.background()
    remove_shadow(shp)
    return shp

def parse_marked(text):
    """==text== 形式を解析し [(text, is_highlighted), ...] を返す"""
    parts = re.split(r'==(.+?)==', text)
    result = []
    for i, part in enumerate(parts):
        if part:
            result.append((part, i % 2 == 1))
    return result

def apply_highlight(run, color_hex="F4F400"):
    """ランの <a:rPr> に <a:highlight> を追加してラインマーカーを適用"""
    r_elem = run._r
    rPr = r_elem.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.Element(qn('a:rPr'))
        r_elem.insert(0, rPr)
    for el in list(rPr):
        if el.tag == qn('a:highlight'):
            rPr.remove(el)
    highlight = etree.SubElement(rPr, qn('a:highlight'))
    etree.SubElement(highlight, qn('a:srgbClr')).set('val', color_hex)

def _run_with_marks(p, text, size, bold=False, color=MAIN2):
    """==text== マーカーを解析し、ハイライト付きランを段落に追加する"""
    for content, highlighted in parse_marked(text):
        r = p.add_run()
        r.text = content
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if highlighted:
            apply_highlight(r)

def _set_para_fmt(p, lsp=150):
    """行間・前後スペース0 を設定"""
    pPr = p._p.get_or_add_pPr()
    for tag in (qn('a:lnSpc'), qn('a:spcBef'), qn('a:spcAft')):
        for el in pPr.findall(tag):
            pPr.remove(el)
    lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
    etree.SubElement(lnSpc, qn('a:spcPct')).set('val', str(lsp * 1000))
    spcB = etree.SubElement(pPr, qn('a:spcBef'))
    etree.SubElement(spcB, qn('a:spcPts')).set('val', '0')
    spcA = etree.SubElement(pPr, qn('a:spcAft'))
    etree.SubElement(spcA, qn('a:spcPts')).set('val', '0')

def _run(p, text, size, bold=False, color=MAIN2):
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color

def add_simple_tb(slide, l, t, w, h, text,
                  size=12, bold=False, color=MAIN1,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.TOP,
                  pl=0.25, pt_=0.25, pr=0.25, pb=0.25, wrap=True):
    """単一段落テキストボックス"""
    tx = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tx.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = v_anchor
    tf.margin_left = Cm(pl);  tf.margin_right  = Cm(pr)
    tf.margin_top  = Cm(pt_); tf.margin_bottom = Cm(pb)
    p = tf.paragraphs[0]
    p.alignment = align
    _set_para_fmt(p, lsp=100)
    _run(p, text, size, bold, color)
    return tx

def add_body_tb(slide, l, t, w, h, text,
                size=13, color=MAIN2,
                pl=0, pr=0, pt_=0.5, pb=0.5):
    """本文テキストボックス（1.5行間・空行除去）"""
    lines = [ln for ln in text.split('\n') if ln.strip()]
    tx = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Cm(pl);  tf.margin_right  = Cm(pr)
    tf.margin_top  = Cm(pt_); tf.margin_bottom = Cm(pb)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        _set_para_fmt(p, lsp=150)
        _run_with_marks(p, line, size, False, color)
    return tx

def add_headed_tb(slide, l, t, w, h, heading, body_text,
                  hsize=13, bsize=11,
                  pl=1.0, pr=1.0, pt_=0.25, pb=0.25):
    """見出し＋本文（1.5行間）テキストボックス"""
    lines = [ln for ln in body_text.split('\n') if ln.strip()]
    tx = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Cm(pl);  tf.margin_right  = Cm(pr)
    tf.margin_top  = Cm(pt_); tf.margin_bottom = Cm(pb)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p0, lsp=150)
    _run(p0, heading, hsize, bold=True, color=ACCENT)
    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        _set_para_fmt(p, lsp=150)
        _run_with_marks(p, line, bsize, False, MAIN2)
    return tx

def add_combined_case_tb(slide, l, t, w, h, overview, flow,
                          hsize=13, bsize=11,
                          pl=1.0, pr=1.0, pt_=0.25, pb=0.25):
    """キャンペーン概要＋7pt空白行＋キャンペーンフロー を1つのTBに統合（1.5行間）"""
    tx = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Cm(pl);  tf.margin_right  = Cm(pr)
    tf.margin_top  = Cm(pt_); tf.margin_bottom = Cm(pb)

    def para(first=False):
        return tf.paragraphs[0] if first else tf.add_paragraph()

    # 「キャンペーン概要」見出し（アクセントカラー）
    p = para(first=True)
    p.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p, lsp=150)
    _run(p, "キャンペーン概要", hsize, bold=True, color=ACCENT)

    # 概要本文（1段落・自然折り返し、==mark== ハイライト対応）
    p = para()
    p.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p, lsp=150)
    _run_with_marks(p, overview, bsize, False, MAIN2)

    # 7pt 空白行（スペース文字で行高を確保）
    p = para()
    p.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p, lsp=150)
    _run(p, " ", 7, False, MAIN2)

    # 「キャンペーンフロー」見出し（アクセントカラー）
    p = para()
    p.alignment = PP_ALIGN.LEFT
    _set_para_fmt(p, lsp=150)
    _run(p, "キャンペーンフロー", hsize, bold=True, color=ACCENT)

    # フロー各ステップ（1ステップ＝1段落、==mark== ハイライト対応）
    for line in [ln for ln in flow.split('\n') if ln.strip()]:
        p = para()
        p.alignment = PP_ALIGN.LEFT
        _set_para_fmt(p, lsp=150)
        _run_with_marks(p, line, bsize, False, MAIN2)

    return tx

def est_h_combined(overview, flow, w_cm,
                   hsize=13, bsize=11, lsp=1.5,
                   ml=1.0, mr=1.0, mt=0.25, mb=0.25, buf=1.12):
    """概要＋スペーサー＋フロー 統合TBの高さを推定（cm）"""
    avail_pt = (w_cm - ml - mr) * 28.35
    cpline = max(1.0, avail_pt / (bsize * 0.85))

    def lh(pt):
        return pt * lsp * 2.54 / 72

    h = mt + mb
    h += lh(hsize)                                              # 概要見出し
    h += math.ceil(max(1, len(overview)) / cpline) * lh(bsize) # 概要本文
    h += lh(7)                                                  # 7pt空白行
    h += lh(hsize)                                              # フロー見出し
    for line in [ln for ln in flow.split('\n') if ln.strip()]:
        h += math.ceil(max(1, len(line)) / cpline) * lh(bsize) # フロー各行
    return h * buf

def est_h(heading, body_text, hsize, bsize, lsp, w_cm,
          ml=1.0, mr=1.0, mt=0.25, mb=0.25, buf=1.12):
    """見出し+本文TBの高さを推定（cm）"""
    avail_pt = (w_cm - ml - mr) * 28.35
    cpline = max(1.0, avail_pt / (bsize * 0.85))

    def lh(pt):
        return pt * lsp * 2.54 / 72

    h = mt + mb + lh(hsize)
    for line in [ln for ln in body_text.split('\n') if ln.strip()]:
        h += math.ceil(max(1, len(line)) / cpline) * lh(bsize)
    return h * buf

def _disp_width(s):
    """文字列の表示幅を返す（全角=1.0・半角=0.5）。日本語の折り返し推定に使う"""
    w = 0.0
    for ch in s:
        if unicodedata.east_asian_width(ch) in ('F', 'W', 'A'):
            w += 1.0
        else:
            w += 0.5
    return w

def est_h_body(text, size=13, lsp=1.5, w_cm=23.4,
               pt_=0.5, pb=0.5, buf=1.08):
    """本文TBの高さを推定（cm）。一般ページのテキスト量に合わせた動的サイズ計算に使う。
    全角=1.0・半角=0.5の表示幅で1行あたりの収容文字数を計算し、折り返し行数を精度よく推定する"""
    avail_pt = w_cm * 28.35
    cpl = max(1.0, avail_pt / size)   # 1行に収まる全角換算文字数

    def lh(pt):
        return pt * lsp * 2.54 / 72

    lines = [ln for ln in text.split('\n') if ln.strip()]
    h = pt_ + pb
    for line in lines:
        rows = max(1, math.ceil(_disp_width(line) / cpl))
        h += rows * lh(size)
    return h * buf

def add_logo(slide, white=True):
    path = LOGO_W if white else LOGO_B
    if os.path.exists(path):
        slide.shapes.add_picture(path, Cm(23.13), Cm(0.82), Cm(1.72), Cm(0.57))

def add_footer(slide, pn, on_accent=False):
    lc = BASE if on_accent else LINE
    tc = BASE if on_accent else CR
    add_rect(slide, 0.59, 13.58, 24.23, 0.03, lc)
    add_simple_tb(slide, 22.73, 13.58, 2.12, 0.71, str(pn),
                  size=10, color=tc, align=PP_ALIGN.RIGHT,
                  v_anchor=MSO_ANCHOR.MIDDLE)
    add_simple_tb(slide, 0.75, 13.52, 7.55, 0.81, "©NONAME Produce Inc.",
                  size=7, color=tc, v_anchor=MSO_ANCHOR.MIDDLE)


# ─── 実績データ表 ─────────────────────────────────────

def add_results_table(slide, rows):
    """
    事例2枚目スライドの実績データ表を追加（スライド右側）
    rows: list of (action_str, count_int)
    例: [("いいね", 1049984), ("リポスト", 301031), ...]
    """
    n_total = len(rows) + 1  # ヘッダー行 + データ行

    tbl_w = Inches(4.61)
    tbl_h = Inches(3.72)
    tbl_x = Inches(5.0)
    tbl_y = Inches(1.41)

    tbl_shape = slide.shapes.add_table(n_total, 2, tbl_x, tbl_y, tbl_w, tbl_h)
    tbl = tbl_shape.table

    # 列幅（アクション40% / 合計件数60%）
    tbl.columns[0].width = Inches(4.61 * 0.40)
    tbl.columns[1].width = Inches(4.61 * 0.60)

    # 行高（均等分割）
    row_h = Inches(3.72 / n_total)
    for row in tbl.rows:
        row.height = row_h

    LGRAY = RGBColor(0xEE, 0xEE, 0xEE)   # ヘッダー以外の WHITE/BLACK/ACCENT はグローバル定数を使う

    # 行数に応じてパディングを調整（基準5行。6行以上は縮小）
    pad_factor = min(1.0, 5 / max(len(rows), 5))
    pad_tb = Cm(0.29 * pad_factor)
    pad_lr = Cm(0.42 * pad_factor)

    def _cell_fmt(cell, text, bg, fg, align):
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        # ① 枠線を #FCFCFC に設定（XML順序上、fillより先に追加）
        for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
            for old in tcPr.findall(qn(tag)):
                tcPr.remove(old)
            ln = etree.SubElement(tcPr, qn(tag))
            ln.set('w', '9525')
            ln.set('cap', 'flat')
            ln.set('cmpd', 'sng')
            sf = etree.SubElement(ln, qn('a:solidFill'))
            etree.SubElement(sf, qn('a:srgbClr')).set('val', 'FCFCFC')

        # ② 垂直中央揃え（anchorを直接指定）
        tcPr.set('anchor', 'ctr')

        # ③ 背景色（fillは枠線の後に追加される）
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

        # ④ マージン
        cell.margin_top    = pad_tb
        cell.margin_bottom = pad_tb
        cell.margin_left   = pad_lr
        cell.margin_right  = pad_lr

        # ⑤ テキスト設定
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = align
        _set_para_fmt(p, lsp=100)
        r = p.add_run()
        r.text = text
        r.font.name  = "Arial"
        r.font.size  = Pt(11)
        r.font.bold  = True
        r.font.color.rgb = fg

    # ヘッダー行（アクセントカラー背景・白文字）
    _cell_fmt(tbl.cell(0, 0), "アクション", ACCENT, WHITE, PP_ALIGN.LEFT)
    _cell_fmt(tbl.cell(0, 1), "合計件数",   ACCENT, WHITE, PP_ALIGN.RIGHT)

    # データ行（薄グレー背景・黒文字）
    for i, (action, count) in enumerate(rows):
        count_str = f"{int(count):,}件"
        _cell_fmt(tbl.cell(i + 1, 0), action,    LGRAY, BLACK, PP_ALIGN.LEFT)
        _cell_fmt(tbl.cell(i + 1, 1), count_str, LGRAY, BLACK, PP_ALIGN.RIGHT)


# ─── ビジュアル：比較表 ───────────────────────────────

def _tbl_cell(cell, text, bg, fg, align,
              bold=True, size=11, pt_=0.1, pb=0.1, pl=0.25, pr=0.25, line_w='19050'):
    """表セルの共通書式（枠線#FCFCFC・2px・垂直中央・背景色・テキスト）"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # 枠線を #FCFCFC・2px に（fillより先に追加する必要がある。9525 EMU = 1px）
    for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set('w', line_w); ln.set('cap', 'flat'); ln.set('cmpd', 'sng')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        etree.SubElement(sf, qn('a:srgbClr')).set('val', 'FCFCFC')
    tcPr.set('anchor', 'ctr')
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.margin_top = Cm(pt_); cell.margin_bottom = Cm(pb)
    cell.margin_left = Cm(pl); cell.margin_right = Cm(pr)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _set_para_fmt(p, lsp=100)
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg


def _tint(color, ratio):
    """color を白(#FCFCFC)に近づけた薄い色を返す（ratio=濃さ 0〜1）"""
    base = 0xFC
    return RGBColor(
        round(base * (1 - ratio) + color[0] * ratio),
        round(base * (1 - ratio) + color[1] * ratio),
        round(base * (1 - ratio) + color[2] * ratio))


_SYMBOLS = ('◯', '○', '〇', '✓', '✔', '×', '✕', '✗', '△', '▲', '-', '―', 'ー')

def _is_symbol(s):
    """機能比較表のセル値が ◯✕△ などの記号1つかを判定"""
    return s.strip() in _SYMBOLS


# 表の共通カラー
HEADER_GREY = RGBColor(0x92, 0x96, 0x99)  # ヘッダー（非強調列）
LABEL_BG    = RGBColor(0xE4, 0xE4, 0xE4)  # 項目名列（左列）
CELL_N      = RGBColor(0xF2, 0xF2, 0xF2)  # 非強調の通常セル
WHITE = RGBColor(0xFC, 0xFC, 0xFC)
BLACK = RGBColor(0x11, 0x11, 0x11)


def add_comparison_table(slide, x, y, w, h, header, rows, emph_idx=None, symbol_mode=False):
    """比較表を (x,y,w,h)[cm] の領域に描画する（N列対応）。
    header: [対象1, 対象2, ...] / rows: [(軸名, [値1, 値2, ...]), ...]
    emph_idx: 強調する対象列のインデックス（0始まり）。None なら最終列を強調
    symbol_mode: True なら ◯✕△ のセルを大きめ太字で表示（機能比較表）"""
    n_obj  = len(header)
    n_cols = n_obj + 1   # 項目名列 + 対象列
    n_rows = len(rows) + 1

    if emph_idx is None:
        emph_idx = n_obj - 1   # デフォルト：最終列を強調

    accent_tint = _tint(ACCENT, 0.12)   # 強調列のデータセル背景（薄アクセント）

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False

    # 列幅：項目名列 28%、対象列を均等
    label_w = w * 0.28
    obj_w   = (w - label_w) / n_obj
    tbl.columns[0].width = Cm(label_w)
    for c in range(n_obj):
        tbl.columns[c + 1].width = Cm(obj_w)

    rh = Cm(h / n_rows)
    for r in tbl.rows:
        r.height = rh

    # ヘッダー行：左上はグレー、対象列は強調列のみアクセント・他はグレー
    _tbl_cell(tbl.cell(0, 0), "", HEADER_GREY, WHITE, PP_ALIGN.LEFT)
    for j, name in enumerate(header):
        bg = ACCENT if j == emph_idx else HEADER_GREY
        _tbl_cell(tbl.cell(0, j + 1), name, bg, WHITE, PP_ALIGN.CENTER)

    # データ行：項目名=#E4E4E4(太字)、通常セル=#F2F2F2(強調列は薄アクセント)
    for i, (label, vals) in enumerate(rows):
        _tbl_cell(tbl.cell(i + 1, 0), label, LABEL_BG, BLACK, PP_ALIGN.LEFT, bold=True)
        for j in range(n_obj):
            v  = vals[j] if j < len(vals) else ""
            bg = accent_tint if j == emph_idx else CELL_N
            if symbol_mode and _is_symbol(v):
                _tbl_cell(tbl.cell(i + 1, j + 1), v, bg, BLACK,
                          PP_ALIGN.CENTER, bold=True, size=16)
            else:
                _tbl_cell(tbl.cell(i + 1, j + 1), v, bg, BLACK,
                          PP_ALIGN.CENTER, bold=False)


def add_correspondence_table(slide, x, y, w, h, header, rows, emph_idx=None):
    """対応表を (x,y,w,h)[cm] の領域に描画する（項目名列なし）。
    header: [左列名, 右列名] / rows: [[左値, 右値], ...]
    emph_idx: 強調する列インデックス（0始まり）。None なら強調なし"""
    n_cols = len(header)
    n_rows = len(rows) + 1

    accent_tint = _tint(ACCENT, 0.12)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False

    col_w = w / n_cols
    for c in range(n_cols):
        tbl.columns[c].width = Cm(col_w)
    rh = Cm(h / n_rows)
    for r in tbl.rows:
        r.height = rh

    # ヘッダー行：強調列のみアクセント・他はグレー
    for j, name in enumerate(header):
        bg = ACCENT if (emph_idx is not None and j == emph_idx) else HEADER_GREY
        _tbl_cell(tbl.cell(0, j), name, bg, WHITE, PP_ALIGN.CENTER)

    # データ行：通常セル=#F2F2F2（強調列は薄アクセント）
    for i, vals in enumerate(rows):
        for j in range(n_cols):
            v  = vals[j] if j < len(vals) else ""
            bg = accent_tint if (emph_idx is not None and j == emph_idx) else CELL_N
            _tbl_cell(tbl.cell(i + 1, j), v, bg, BLACK, PP_ALIGN.LEFT, bold=False)


# ─── ビジュアル：グラフ（matplotlibで画像化） ───────────────

# 日本語フォントを設定（Windows標準のMeiryo）
for _f in ('Meiryo', 'Yu Gothic', 'MS Gothic'):
    if any(ft.name == _f for ft in font_manager.fontManager.ttflist):
        plt.rcParams['font.family'] = _f
        break
plt.rcParams['axes.unicode_minus'] = False

def _hex(color):
    """RGBColor → matplotlib用 '#RRGGBB' 文字列"""
    return '#%02X%02X%02X' % (color[0], color[1], color[2])

CM_PER_INCH = 2.54

def _save_fig_to_slide(fig, slide, x, y, w, h):
    """matplotlib figure を一時PNG経由でスライドのビジュアル領域に貼る。
    figsize を領域の実寸(cm→inch)で作っている前提で、bbox_inchesを使わず
    そのまま領域(x,y,w,h)[cm]に貼ることで歪みを防ぐ。"""
    tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
    tmp_path = tmp.name
    tmp.close()
    fig.savefig(tmp_path, dpi=200, transparent=True)   # bbox_inches を指定しない
    plt.close(fig)
    slide.shapes.add_picture(tmp_path, Cm(x), Cm(y), Cm(w), Cm(h))
    try:
        os.remove(tmp_path)
    except OSError:
        pass


def add_bar_chart(slide, x, y, w, h, title, data, emph=None):
    """棒グラフを (x,y,w,h)[cm] の領域に描画する。
    data: [(ラベル, 値), ...] / emph: 強調する棒のインデックスのリスト（空/Noneなら全て同色）"""
    labels = [d[0] for d in data]
    values = [d[1] for d in data]
    n = len(values)
    emph = emph or []

    # 値の開きが大きい（最大/最小が30倍以上）なら対数スケール。小さい棒が潰れるのを防ぐ
    nz = [v for v in values if v > 0]
    use_log = bool(nz) and (max(nz) / min(nz) >= 30)

    # figsize = ビジュアル領域の実寸（インチ）→ 歪みなし
    fig, ax = plt.subplots(figsize=(w / CM_PER_INCH, h / CM_PER_INCH))
    # タイトルを下に置くスペースを確保しつつ余白を最小化
    fig.subplots_adjust(top=0.95, bottom=0.22, left=0.04, right=0.97)

    accent = _hex(ACCENT)
    grey   = _hex(HEADER_GREY)
    if emph:
        colors = [accent if i in emph else grey for i in range(n)]
    else:
        colors = [accent] * n

    bars = ax.bar(labels, values, color=colors, width=0.6, zorder=3)

    # 各棒の上に値ラベル
    for b, v in zip(bars, values):
        ax.text(b.get_x() + b.get_width() / 2, b.get_height(),
                f"{v:g}", ha='center', va='bottom', fontsize=11,
                color=_hex(MAIN1), fontweight='bold')

    # 軸の装飾を最小化（罫線なし・枠なし・Y軸目盛なし）
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.spines['bottom'].set_color(_hex(LINE))
    xfs = 11 if n <= 4 else 9   # 本数が多いとX軸ラベルを小さく
    ax.tick_params(axis='x', length=0, labelsize=xfs, colors=_hex(MAIN2))
    ax.tick_params(axis='y', length=0, labelleft=False)
    if use_log:
        ax.set_yscale('log')
        ax.set_ylim(min(nz) * 0.6, max(nz) * 2.2)
    else:
        ax.set_ylim(0, max(values) * 1.18)
    ax.margins(x=0.02)

    # タイトルはグラフの下に配置（X軸ラベルとして）
    if title:
        ax.set_xlabel(title, fontsize=12, color=_hex(MAIN1),
                      fontweight='bold', labelpad=12)

    _save_fig_to_slide(fig, slide, x, y, w, h)


# ─── ビジュアル：フロー図（ネイティブ図形） ─────────────────

def add_flow_diagram(slide, x, y, w, h, steps):
    """フロー図を (x,y,w,h)[cm] の領域に描画する（シェブロン型・段階的な色）。
    steps: [(タイトル, 説明), ...] 説明が空なら番号＋タイトルのみのシンプル表示。
    説明があれば STEP番号＋タイトル＋説明 のリッチ表示になる。"""
    n = len(steps)
    if n == 0:
        return
    has_desc = any(d for (_, d) in steps)

    overlap = 0.45                                   # シェブロンの重なり（矢羽が次に食い込む）
    cw = (w + (n - 1) * overlap) / n                 # シェブロン1個の幅
    ch_h = min(h, 6.2)                               # シェブロン高さ（領域内で上下中央）
    cy = y + (h - ch_h) / 2
    r_min, r_max = 0.06, 0.24                         # 段階的な濃さ（左薄→右濃）
    inset = ch_h * 0.26                              # 矢羽の水平幅ぶんテキストを内側へ

    for i, (title, desc) in enumerate(steps):
        cx = x + i * (cw - overlap)
        ratio = r_min + (r_max - r_min) * (i / max(n - 1, 1))
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON, Cm(cx), Cm(cy), Cm(cw), Cm(ch_h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _tint(ACCENT, ratio)
        shp.line.fill.background()
        remove_shadow(shp)

        # テキストはシェブロンに重ねた別TBで配置（矢羽形状を避けて内側に置く）
        tx = cx + inset
        tw = cw - inset * 1.25
        tb = slide.shapes.add_textbox(Cm(tx), Cm(cy), Cm(tw), Cm(ch_h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Cm(0.1); tf.margin_right = Cm(0.1)
        tf.margin_top = Cm(0.1); tf.margin_bottom = Cm(0.1)

        # STEP番号
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        _set_para_fmt(p, lsp=100)
        _run(p, f"STEP {str(i + 1).zfill(2)}", 10, bold=True, color=ACCENT)

        # タイトル
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        _set_para_fmt(p2, lsp=110)
        _run(p2, title, 12 if has_desc else 13, bold=True, color=MAIN1)

        # 説明（あれば）
        if desc:
            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.LEFT
            _set_para_fmt(p3, lsp=115)
            _run(p3, desc, 9, bold=False, color=MAIN2)


# ─── スライドビルダー ──────────────────────────────────

def build_cover(prs, pn):
    # 表紙は白紙（ページ番号・©・下線なし）。背景のみ
    sl = blank_slide(prs)
    set_bg(sl, BASE)

def build_toc(prs, chapters, pn):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_rect(sl, 12.5, 0, 12.95, 14.29, ACCENT)
    add_logo(sl, white=True)
    add_rect(sl, 0.59, 13.58, 24.23, 0.03, LINE)
    add_simple_tb(sl, 0.75, 13.52, 7.55, 0.81, "©NONAME Produce Inc.",
                  size=7, color=CR, v_anchor=MSO_ANCHOR.MIDDLE)
    add_simple_tb(sl, 22.73, 13.58, 2.12, 0.71, str(pn),
                  size=10, color=BASE, align=PP_ALIGN.RIGHT, v_anchor=MSO_ANCHOR.MIDDLE)
    add_simple_tb(sl, 1, 5.87, 1.72, 1.11, "00.",
                  size=14, color=MAIN1, v_anchor=MSO_ANCHOR.MIDDLE)
    add_simple_tb(sl, 1.79, 5.83, 6.13, 1.54, "目次",
                  size=24, bold=True, color=MAIN1, v_anchor=MSO_ANCHOR.MIDDLE)

    # 章リスト TB（幅3.84in × 高4.5in、X=5.27in、スライド上下中央）
    tb_w = Inches(3.84)
    tb_h = Inches(4.5)
    tb_x = Inches(5.27)
    tb_y = Inches((14.29 / 2.54 - 4.5) / 2)

    lines = []
    for i, ch in enumerate(chapters):
        lines.append(f"{str(i+1).zfill(2)}.{ch}")
        if i < len(chapters) - 1:
            lines.append("")

    tx = sl.shapes.add_textbox(tb_x, tb_y, tb_w, tb_h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Cm(0.25); tf.margin_right  = Cm(0.25)
    tf.margin_top  = Cm(0.25); tf.margin_bottom = Cm(0.25)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if line.strip():
            _set_para_fmt(p, lsp=115)
            _run(p, line, 14, bold=True, color=BASE)
        else:
            _set_para_fmt(p, lsp=115)
            _run(p, "", 7, color=BASE)

def build_chapter(prs, ch_num, title, pn):
    """章題ページ。タイトルは左上揃え。"""
    sl = blank_slide(prs)
    set_bg(sl, ACCENT)
    add_logo(sl, white=True)
    add_footer(sl, pn, on_accent=True)
    add_simple_tb(sl, 1, 5.87, 1.72, 1.11, f"{str(ch_num).zfill(2)}.",
                  size=14, color=BASE, v_anchor=MSO_ANCHOR.MIDDLE)
    # タイトル：左揃え・上揃え（spec: 左上）
    add_simple_tb(sl, 1.79, 5.83, 19.68, 2.5, title,
                  size=24, bold=True, color=BASE,
                  v_anchor=MSO_ANCHOR.TOP, wrap=True,
                  pl=0.25, pt_=0.25, pr=0.25, pb=0.25)

def build_general(prs, title, body, pn, visual=None):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_rect(sl, 0, 1.08, 1, 0.03, BAR)
    add_logo(sl, white=False)
    add_footer(sl, pn)
    add_simple_tb(sl, 1, 0.69, 21.58, 0.81, title,
                  size=19, bold=True, color=ACCENT,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)
    body_h = est_h_body(body)
    add_body_tb(sl, 1, 1.51, 23.4, body_h, body, size=13, color=MAIN2)

    # ── ビジュアル領域 ──
    # 上辺: 本文TBの下辺 / 下辺: 下部ライン(Y13.58)の0.27インチ上 / 横幅: 本文と同じ
    vis_x      = 1.0
    vis_w      = 23.4
    vis_top    = 1.51 + body_h
    vis_bottom = 13.58 - 0.27 * 2.54   # 0.27インチ = 0.6858cm
    vis_h      = vis_bottom - vis_top
    if vis_h <= 0.5:
        return

    vtype = visual.get('type') if visual else None
    if vtype == '比較表':
        add_comparison_table(sl, vis_x, vis_top, vis_w, vis_h,
                             visual.get('header'), visual.get('rows'),
                             emph_idx=visual.get('emph_idx'))
    elif vtype == '機能比較表':
        add_comparison_table(sl, vis_x, vis_top, vis_w, vis_h,
                             visual.get('header'), visual.get('rows'),
                             emph_idx=visual.get('emph_idx'), symbol_mode=True)
    elif vtype == '対応表':
        add_correspondence_table(sl, vis_x, vis_top, vis_w, vis_h,
                             visual.get('header'), visual.get('rows'),
                             emph_idx=visual.get('emph_idx'))
    elif vtype == '棒グラフ':
        add_bar_chart(sl, vis_x, vis_top, vis_w, vis_h,
                      visual.get('title'), visual.get('data'),
                      emph=visual.get('emph'))
    elif vtype == 'フロー図':
        add_flow_diagram(sl, vis_x, vis_top, vis_w, vis_h, visual.get('steps'))
    else:
        # 未実装の種類 / なし はプレースホルダ表示
        add_rect(sl, vis_x, vis_top, vis_w, vis_h, RGBColor(0xEC, 0xEC, 0xE6))
        label = f"ビジュアル領域（{vtype}）" if vtype else "ビジュアル領域"
        add_simple_tb(sl, vis_x, vis_top, vis_w, vis_h, label,
                      size=14, color=MAIN2, align=PP_ALIGN.CENTER,
                      v_anchor=MSO_ANCHOR.MIDDLE)

def build_case1(prs, case_num, name, company, overview, flow, pn):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_rect(sl, 0, 0, 25.4, 3.01, ACCENT)
    add_rect(sl, 0, 1.08, 1, 0.03, BASE)
    add_logo(sl, white=True)
    add_footer(sl, pn)

    add_simple_tb(sl, 1, 0.69, 21.58, 0.81,
                  f"{str(case_num).zfill(2)}.{name}",
                  size=19, bold=True, color=BASE,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)
    add_simple_tb(sl, 1, 1.76, 21.58, 0.56, f"by {company}",
                  size=13, color=BASE,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)

    tb_w = 12.7
    tb_h = est_h_combined(overview, flow, tb_w)
    add_combined_case_tb(sl, 0, 3.3, tb_w, tb_h, overview, flow)

def build_case2(prs, case_num, name, company, point, pn, results=None):
    sl = blank_slide(prs)
    set_bg(sl, BASE)
    add_rect(sl, 0, 0, 25.4, 3.01, ACCENT)
    add_rect(sl, 0, 1.08, 1, 0.03, BASE)
    add_logo(sl, white=True)
    add_footer(sl, pn)

    add_simple_tb(sl, 1, 0.69, 21.58, 0.81,
                  f"{str(case_num).zfill(2)}.{name}",
                  size=19, bold=True, color=BASE,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)
    add_simple_tb(sl, 1, 1.76, 21.58, 0.56, f"by {company}",
                  size=13, color=BASE,
                  align=PP_ALIGN.LEFT, v_anchor=MSO_ANCHOR.MIDDLE,
                  pl=0.5, pt_=0.25, pr=0.25, pb=0.25)

    pt_h = est_h("キャンペーンのポイント", point, 13, 11, 1.5, 12.7)
    add_headed_tb(sl, 0, 3.3, 12.7, pt_h,
                  "キャンペーンのポイント", point, hsize=13, bsize=11)
    if results:
        add_results_table(sl, results)


# ─── PPTX 差し込み（PowerPoint COM を使用）────────────────

def _com_insert_service_slides(out_path, service_map, ch_positions):
    """
    PowerPoint COM を使ってサービスファイルを章タイトルの直後に差し込む。
    テーマ・フォント・カラーをソース通りに保持する唯一の確実な方法。

    service_map:  {chapter_num: src_file_path}
    ch_positions: {chapter_num: 1-indexed slide number in target}
    """
    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True   # False だと一部操作が失敗するため True で起動

    skipped = []
    try:
        tgt = ppt.Presentations.Open(
            os.path.abspath(out_path),
            False,   # ReadOnly
            False,   # Untitled
            False)   # WithWindow（バックグラウンド）

        offset = 0  # 先の挿入でスライド番号がずれる分を補正

        for ch in sorted(ch_positions.keys()):
            src = service_map.get(ch, "")
            insert_after = ch_positions[ch] + offset

            if not src or not os.path.exists(src):
                skipped.append(f"第{ch}章（{os.path.basename(src or '未設定')} が見つかりません）")
                continue

            abs_src = os.path.abspath(src)

            # スライド枚数を確認
            tmp = ppt.Presentations.Open(abs_src, True, False, False)
            n = tmp.Slides.Count
            tmp.Close()

            # 差し込み（ソース書式を保持）
            tgt.Slides.InsertFromFile(abs_src, insert_after)
            offset += n
            print(f"  INSERT: 第{ch}章 → {n}枚 after slide {insert_after}")

        tgt.Save()
        tgt.Close()

    finally:
        ppt.Quit()
        pythoncom.CoUninitialize()

    return skipped


# ─── Markdownパーサー ──────────────────────────────────

def _parse_results_field(text):
    """'いいね:1049984 / リポスト:301031 / ...' → [(str, int), ...] or None"""
    rows = []
    for item in text.split(' / '):
        item = item.strip()
        if ':' not in item:
            continue
        action, _, raw_count = item.rpartition(':')
        try:
            rows.append((action.strip(), int(raw_count.strip().replace(',', ''))))
        except ValueError:
            pass
    return rows if rows else None


def _parse_flow_text(text):
    """'1. step1 2. step2 ...' → '1. step1\n2. step2\n...'"""
    text = text.strip()
    return re.sub(r'\s+(\d+\.\s)', r'\n\1', text)


def _parse_page_fields(block):
    """ページブロック内のフィールド行をパースして dict を返す"""
    fields = {}
    cur_key = None
    cur_val = []

    for line in block.split('\n'):
        stripped = line.strip()
        # 【図】【ビジュアル】行はスキップ
        if re.match(r'^-\s*【', stripped):
            if cur_key:
                fields[cur_key] = ' '.join(cur_val).strip()
                cur_key, cur_val = None, []
            continue
        # フィールド行（- キー：値）
        m = re.match(r'^-\s*([^：\[]+?)：(.*)', stripped)
        if m:
            if cur_key:
                fields[cur_key] = ' '.join(cur_val).strip()
            cur_key = m.group(1).strip()
            cur_val = [m.group(2).strip()]
        elif stripped and cur_key:
            cur_val.append(stripped)

    if cur_key:
        fields[cur_key] = ' '.join(cur_val).strip()

    return fields


def _parse_visual(block):
    """ページブロックから 【ビジュアル】 を抽出して dict を返す。なし/未対応なら None。
    比較表のみ実装済み。他種類は {'type': 種類名} を返しプレースホルダ表示に使う。"""
    lines = block.split('\n')
    vtype = None
    data_lines = []
    capturing = False

    for line in lines:
        m = re.match(r'^-\s*【ビジュアル】\s*[:：]\s*(.+)', line.strip())
        if m:
            vtype = m.group(1).strip()
            capturing = True
            continue
        if capturing:
            if line.startswith(' ') or line.startswith('　'):
                if line.strip():
                    data_lines.append(line.strip())
            elif line.strip().startswith('-') or line.strip().startswith('#'):
                break
            elif not line.strip():
                continue
            else:
                break

    if not vtype or vtype.startswith('なし'):
        return None

    def _split_emph(parts):
        """各列名から ★/☆ を検出し、(クリーンな列名リスト, 強調インデックス) を返す"""
        clean = []
        emph = None
        for idx, p in enumerate(parts):
            if p.startswith('★') or p.startswith('☆'):
                emph = idx
                p = p.lstrip('★☆').strip()
            clean.append(p)
        return clean, emph

    if vtype.startswith('比較表') or vtype.startswith('機能比較表'):
        header = None
        emph_idx = None
        rows = []
        for dl in data_lines:
            key, sep, val = dl.partition('：')
            if not sep:
                key, sep, val = dl.partition(':')
            key = key.strip()
            parts = [p.strip() for p in val.split('|')]
            if key == '対象':
                header, emph_idx = _split_emph(parts)
            else:
                rows.append((key, parts))
        ttype = '機能比較表' if vtype.startswith('機能比較表') else '比較表'
        if header and rows:
            return {'type': ttype, 'header': header,
                    'rows': rows, 'emph_idx': emph_idx}
        return None

    if vtype.startswith('対応表'):
        header = None
        emph_idx = None
        rows = []
        for dl in data_lines:
            if dl.startswith('列：') or dl.startswith('列:'):
                _, sep, val = dl.partition('：')
                if not sep:
                    _, sep, val = dl.partition(':')
                parts = [p.strip() for p in val.split('|')]
                header, emph_idx = _split_emph(parts)
            else:
                rows.append([p.strip() for p in dl.split('|')])
        if header and rows:
            return {'type': '対応表', 'header': header,
                    'rows': rows, 'emph_idx': emph_idx}
        return None

    if vtype.startswith('棒グラフ'):
        title = ''
        data = []
        emph = []          # ★を付けた棒のインデックス（複数可）
        for dl in data_lines:
            key, sep, val = dl.partition('：')
            if not sep:
                key, sep, val = dl.partition(':')
            key = key.strip(); val = val.strip()
            if key == 'タイトル':
                title = val
            elif key == 'データ':
                idx = 0
                for item in val.split(' / '):
                    item = item.strip()
                    if not item:
                        continue
                    is_emph = item.startswith('★') or item.startswith('☆')
                    if is_emph:
                        item = item.lstrip('★☆').strip()
                    label, sep2, raw = item.rpartition(':')
                    if not sep2:
                        label, sep2, raw = item.rpartition('：')
                    try:
                        data.append((label.strip(), float(raw.strip().replace(',', ''))))
                        if is_emph:
                            emph.append(idx)
                        idx += 1
                    except ValueError:
                        pass
        if data:
            return {'type': '棒グラフ', 'title': title,
                    'data': data, 'emph': emph}
        return None

    if vtype.startswith('フロー図'):
        step_lines = []   # 「ステップ: タイトル ; 説明」形式（複数行）
        arrow_val  = None  # 「ステップ: A → B → C」形式（1行・説明なし）
        for dl in data_lines:
            key, sep, val = dl.partition('：')
            if not sep:
                key, sep, val = dl.partition(':')
            if key.strip() == 'ステップ':
                val = val.strip()
                if re.search(r'[→➡⇒]', val):
                    arrow_val = val
                else:
                    step_lines.append(val)
        steps = []   # [(title, desc), ...]
        if step_lines:
            for sl_ in step_lines:
                parts = re.split(r'[;；]', sl_, maxsplit=1)
                title = parts[0].strip()
                desc  = parts[1].strip() if len(parts) > 1 else ''
                if title:
                    steps.append((title, desc))
        elif arrow_val:
            for s in re.split(r'[→➡⇒]', arrow_val):
                s = s.strip()
                if s:
                    steps.append((s, ''))
        if steps:
            return {'type': 'フロー図', 'steps': steps}
        return None

    # 未実装の種類はプレースホルダ用に種類名のみ返す
    return {'type': vtype}


def _resolve_insert_file(ch_title):
    """章タイトルから差し込みPPTXのパスを返す。見つからなければ None"""
    if 'その他のご案内' in ch_title:
        p = os.path.join(INSERT_DIR, 'n2p-wp-other.pptx')
        return p if os.path.exists(p) else None
    # INSERT_DIR 内の n2p-wp-*.pptx（other 以外）を名前順で最初に見つかるものを使う
    if os.path.isdir(INSERT_DIR):
        for fname in sorted(os.listdir(INSERT_DIR)):
            if (fname.startswith('n2p-wp-')
                    and fname.endswith('.pptx')
                    and fname != 'n2p-wp-other.pptx'):
                return os.path.join(INSERT_DIR, fname)
    return None


def parse_markdown(md_path):
    """
    WP構成Markdownをパースしてスライドデータを返す。
    戻り値: {'wp_title': str, 'chapters': [str], 'pages': [dict]}
    page dict の type:
      'chapter'  → {ch_num, title}
      'service'  → {ch_num, title}   ← 〇〇のご案内
      'general'  → {title, body}
      'case'     → {title, company, overview, flow, point, results}
    """
    with open(md_path, encoding='utf-8') as f:
        raw = f.read()

    # タイトル候補リスト以降は不要
    raw = re.split(r'\n## タイトル候補', raw)[0]

    # WPタイトル
    m = re.search(r'^# タイトル（WP名）：(.+)$', raw, re.MULTILINE)
    wp_title = m.group(1).strip() if m else 'ホワイトペーパー'

    pages = []
    chapter_titles = []

    for block in re.split(r'\n##\s+', raw)[1:]:
        lines = block.split('\n')
        ch_header = lines[0].strip()

        cm = re.match(r'第(\d+)章：(.+)', ch_header)
        if not cm:
            continue
        ch_num   = int(cm.group(1))
        ch_title = cm.group(2).strip()
        chapter_titles.append(ch_title)

        if ch_title.endswith('のご案内'):
            pages.append({'type': 'service', 'ch_num': ch_num, 'title': ch_title})
            continue

        pages.append({'type': 'chapter', 'ch_num': ch_num, 'title': ch_title})

        for pb in re.split(r'\n###\s+', '\n'.join(lines[1:]))[1:]:
            pg_lines = pb.split('\n')
            pm = re.match(r'[\d]+-[\d]+：(.+)', pg_lines[0].strip())
            if not pm:
                continue
            pg_title = pm.group(1).strip()
            fields   = _parse_page_fields('\n'.join(pg_lines[1:]))

            if 'フロー' in fields:
                pages.append({
                    'type':     'case',
                    'title':    pg_title,
                    'company':  fields.get('企業名', ''),
                    'overview': fields.get('内容', ''),
                    'flow':     _parse_flow_text(fields.get('フロー', '')),
                    'point':    fields.get('ポイント', ''),
                    'results':  _parse_results_field(fields['実績']) if '実績' in fields else None,
                })
            else:
                pages.append({
                    'type':   'general',
                    'title':  pg_title,
                    'body':   fields.get('内容', ''),
                    'visual': _parse_visual(pb),
                })

    return {'wp_title': wp_title, 'chapters': chapter_titles, 'pages': pages}


# ─── メイン ────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("Usage: python generate_slides.py <*-wp-structure.md> [ACCENT_RRGGBB]")
        sys.exit(1)

    md_path = sys.argv[1]

    # アクセントカラーをコマンドライン引数で上書き可能（例: 1665F2）
    global ACCENT
    if len(sys.argv) >= 3:
        h = sys.argv[2].lstrip('#')
        ACCENT = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    data = parse_markdown(md_path)

    # 出力先: Markdownと同じフォルダ・同じプロジェクト名
    md_dir  = os.path.dirname(os.path.abspath(md_path))
    proj    = os.path.basename(md_path).replace('-wp-structure.md', '')
    out     = os.path.join(md_dir, f"{proj}-slides.pptx")

    prs = Presentation()
    prs.slide_width  = Cm(25.4)
    prs.slide_height = Cm(14.29)
    pn = 1

    build_cover(prs, pn); pn += 1
    build_toc(prs, data['chapters'], pn); pn += 1

    case_cnt      = 0
    svc_positions = {}   # {ch_num: 1-indexed slide position}
    svc_files     = {}   # {ch_num: insert_file_path}

    for page in data['pages']:
        t = page['type']

        if t == 'chapter':
            build_chapter(prs, page['ch_num'], page['title'], pn); pn += 1

        elif t == 'service':
            build_chapter(prs, page['ch_num'], page['title'], pn); pn += 1
            svc_positions[page['ch_num']] = len(prs.slides)
            svc_files[page['ch_num']]     = _resolve_insert_file(page['title'])

        elif t == 'general':
            build_general(prs, page['title'], page['body'], pn,
                          visual=page.get('visual')); pn += 1

        elif t == 'case':
            case_cnt += 1
            build_case1(prs, case_cnt, page['title'], page['company'],
                        page['overview'], page['flow'], pn); pn += 1
            build_case2(prs, case_cnt, page['title'], page['company'],
                        page['point'], pn, results=page['results']); pn += 1

    prs.save(out)
    print(f"本文スライド生成完了: {len(prs.slides)}枚 -> {out}")

    if svc_positions:
        skipped = _com_insert_service_slides(out, svc_files, svc_positions)
        print("差し込み完了（PowerPoint COM）")
        for s in skipped:
            print(f"  SKIP: {s}")

if __name__ == "__main__":
    main()
